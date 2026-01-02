import os

# 导入Office文档处理器
try:
    from docx import Document
    DOCX_SUPPORTED = True
except ImportError:
    DOCX_SUPPORTED = False

try:
    from openpyxl import load_workbook
    XLSX_SUPPORTED = True
except ImportError:
    XLSX_SUPPORTED = False

try:
    from pptx import Presentation
    PPTX_SUPPORTED = True
except ImportError:
    PPTX_SUPPORTED = False

try:
    import xlrd
    XLS_SUPPORTED = True
except ImportError:
    XLS_SUPPORTED = False

try:
    import win32com.client
    import pythoncom
    WIN32COM_SUPPORTED = True
except ImportError:
    WIN32COM_SUPPORTED = False

class OfficeProcessor:
    """Office文档处理模块，负责处理Word、Excel、PowerPoint等Office文档"""
    
    def __init__(self):
        pass
    
    def read_docx_content(self, docx_path):
        """读取Word文档内容"""
        if not DOCX_SUPPORTED:
            return "未安装python-docx库，无法读取Word文档"
        try:
            doc = Document(docx_path)
            content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    table_content.append(row_text)
                if table_content:
                    content.append("\n表格:")
                    content.extend(table_content)
            return "\n".join(content)
        except Exception as e:
            return f"读取Word文档失败: {str(e)}"
    
    def read_xlsx_content(self, xlsx_path):
        """读取Excel文档内容"""
        if not XLSX_SUPPORTED:
            return "未安装openpyxl库，无法读取Excel文档"
        try:
            wb = load_workbook(xlsx_path)
            content = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                content.append(f"\n工作表 {sheet_name}")
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        content.append(row_text)
            return "\n".join(content)
        except Exception as e:
            return f"读取Excel文档失败: {str(e)}"
    
    def read_pptx_content(self, pptx_path):
        """读取PowerPoint文档内容"""
        if not PPTX_SUPPORTED:
            return "未安装python-pptx库，无法读取PowerPoint文档"
        try:
            prs = Presentation(pptx_path)
            content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                content.append(f"\n幻灯片{slide_num}:")
                if hasattr(slide.shapes, 'title') and slide.shapes.title:
                    if slide.shapes.title.text.strip():
                        content.append(f"标题: {slide.shapes.title.text}")
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        content.append(f"内容: {shape.text_frame.text}")
            return "\n".join(content)
        except Exception as e:
            return f"读取PowerPoint文档失败: {str(e)}"
    
    def read_doc_content(self, doc_path):
        """读取旧版Word文档内容"""
        if not WIN32COM_SUPPORTED:
            return "未安装pywin32库，无法读取旧版Word文档(.doc)"
        try:
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(doc_path)
            content = doc.Content.Text
            doc.Close()
            word.Quit()
            pythoncom.CoUninitialize()
            return content
        except Exception as e:
            return f"读取旧版Word文档失败: {str(e)}"
    
    def read_xls_content(self, xls_path):
        """读取旧版Excel文档内容"""
        if not XLS_SUPPORTED:
            return "未安装xlrd库，无法读取旧版Excel文档(.xls)"
        try:
            workbook = xlrd.open_workbook(xls_path)
            content = []
            for sheet_name in workbook.sheet_names():
                sheet = workbook.sheet_by_name(sheet_name)
                content.append(f"\n工作表 {sheet_name}")
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_text = " | ".join([str(cell) if cell != "" else "" for cell in row])
                    if row_text.strip():
                        content.append(row_text)
            return "\n".join(content)
        except Exception as e:
            return f"读取旧版Excel文档失败: {str(e)}"
    
    def read_ppt_content(self, ppt_path):
        """读取旧版PowerPoint文档内容"""
        if not WIN32COM_SUPPORTED:
            return "未安装pywin32库，无法读取旧版PowerPoint文档(.ppt)"
        try:
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            prs = powerpoint.Presentations.Open(ppt_path, ReadOnly=True)
            content = []
            for slide_num, slide in enumerate(prs.Slides, 1):
                content.append(f"\n幻灯片{slide_num}:")
                try:
                    title = slide.Shapes.Title.TextFrame.TextRange.Text
                    if title.strip():
                        content.append(f"标题: {title}")
                except Exception:
                    pass
                for shape in slide.Shapes:
                    try:
                        if shape.HasTextFrame:
                            text = shape.TextFrame.TextRange.Text
                            if text.strip():
                                content.append(f"内容: {text}")
                    except Exception:
                        pass
            prs.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
            return "\n".join(content)
        except Exception as e:
            return f"读取旧版PowerPoint文档失败: {str(e)}"
